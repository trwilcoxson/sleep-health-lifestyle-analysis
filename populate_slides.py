"""
Populate the Sleep Health and Lifestyle Analysis slide deck template.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

prs = Presentation('template.pptx')

# ============================================================
# Strip blue circle bullets from ALL slide layouts at source
# ============================================================
from pptx.oxml.ns import qn

for layout in prs.slide_layouts:
    for shape in layout.placeholders:
        if shape.has_text_frame:
            # Find lstStyle in the text body
            txBody = shape._element.find(qn('p:txBody'))
            if txBody is None:
                continue
            lstStyle = txBody.find(qn('a:lstStyle'))
            if lstStyle is None:
                continue
            # For each level, remove bullet chars/colors and add buNone
            for lvl in lstStyle:
                for child in list(lvl):
                    tag = child.tag
                    if any(x in tag for x in ['buChar', 'buClr', 'buSzPts', 'buSzPct', 'buFont', 'buAutoNum']):
                        lvl.remove(child)
                # Add buNone and reset indent/margin
                lvl.append(lvl.makeelement(qn('a:buNone'), {}))
                lvl.set('indent', '0')
                lvl.set('marL', '0')

# Also strip from slide masters
for master in prs.slide_masters:
    for shape in master.placeholders:
        if shape.has_text_frame:
            txBody = shape._element.find(qn('p:txBody'))
            if txBody is None:
                continue
            lstStyle = txBody.find(qn('a:lstStyle'))
            if lstStyle is None:
                continue
            for lvl in lstStyle:
                for child in list(lvl):
                    tag = child.tag
                    if any(x in tag for x in ['buChar', 'buClr', 'buSzPts', 'buSzPct', 'buFont', 'buAutoNum']):
                        lvl.remove(child)
                lvl.append(lvl.makeelement(qn('a:buNone'), {}))
                lvl.set('indent', '0')
                lvl.set('marL', '0')

# ============================================================
# Helper: set paragraph text with formatting
# ============================================================
def clear_and_set_text(text_frame, lines, font_size=14, bold_first=False):
    """Clear text frame and add lines with formatting."""
    from pptx.oxml.ns import qn

    # Clear existing paragraphs
    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
        p = text_frame.paragraphs[i]._p
        p.getparent().remove(p)

    # Clear the first paragraph's existing runs
    first_para = text_frame.paragraphs[0]
    for run in first_para.runs:
        run._r.getparent().remove(run._r)

    for i, line in enumerate(lines):
        if i == 0:
            para = first_para
        else:
            para = text_frame.add_paragraph()

        # Remove template bullet formatting (blue circles from layout)
        pPr = para._p.get_or_add_pPr()
        # Reset indent and margin to prevent bullet space
        pPr.set('indent', '0')
        pPr.set('marL', '0')
        # Remove any existing bullet-related elements
        for child in list(pPr):
            if any(x in child.tag for x in ['buChar', 'buAutoNum', 'buFont', 'buSzPts', 'buSzPct', 'buClr', 'buNone']):
                pPr.remove(child)
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))

        # Handle bold segments with (text, bold) tuples
        if isinstance(line, list):
            for segment in line:
                if isinstance(segment, tuple):
                    text, is_bold = segment
                    run = para.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
                    run.font.bold = is_bold
                else:
                    run = para.add_run()
                    run.text = segment
                    run.font.size = Pt(font_size)
        elif line == "":
            para.space_after = Pt(4)
        else:
            run = para.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            if bold_first and i == 0:
                run.font.bold = True


# ============================================================
# Slide 1: Title Slide - Update subtitle and name
# ============================================================
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '[Subtitle]' in run.text:
                    run.text = run.text.replace('[Subtitle]', 'Descriptive Statistics and Data Exploration')
                if '[Your Name]' in run.text:
                    run.text = run.text.replace('[Your Name]', 'Tim Wilcoxson')

# ============================================================
# Slide 2: Delete "How to Use This Template" slide
# ============================================================
rId = prs.slides._sldIdLst[1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[1]

# After deletion, slide indices shift:
# Slide 0 = Title, Slide 1 = Agenda, Slide 2 = Data Description,
# Slide 3 = Physical Activity, Slide 4 = Daily Steps, Slide 5 = Heart Rate

# ============================================================
# Slide 3 (index 2): Data Description
# ============================================================
slide_data = prs.slides[2]
for shape in slide_data.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '[TODO]' in para.text:
                content = [
                    [("Sleep Duration (hours)", True), (" is an example of a ", False), ("continuous", True), (" variable in the dataset.", False)],
                    "    Values such as 5.9, 6.1, and 7.8 can take any real number within a range.",
                    "",
                    [("Daily Steps", True), (" is an example of an ", False), ("integer", True), (" variable in the dataset.", False)],
                    "    Values such as 3,000, 5,600, and 10,000 are whole numbers.",
                    "",
                    [("Quality of Sleep (scale: 1\u201310)", True), (" is an example of an ", False), ("ordinal categorical", True), (" variable in the dataset.", False)],
                    "    The ratings have a meaningful order from low (1) to high (10).",
                    "",
                    [("Gender", True), (" is an example of a ", False), ("nominal categorical", True), (" variable in the dataset.", False)],
                    "    Categories (Male, Female) have no inherent order or ranking.",
                ]
                clear_and_set_text(shape.text_frame, content, font_size=14)
                break

# ============================================================
# Slide 4 (index 3): Typical Amount of Physical Activity
# ============================================================
slide_pa = prs.slides[3]
for shape in slide_pa.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '[TODO]' in para.text:
                content = [
                    [("Measures of Center", True)],
                    "",
                    [("\u2022  Mean: ", True), ("59.17 minutes per day", False)],
                    [("\u2022  Median: ", True), ("60.00 minutes per day", False)],
                    [("\u2022  Mode: ", True), ("60 minutes per day", False)],
                    "",
                    [("Distribution Shape (Mean vs. Median)", True)],
                    "",
                    "The mean (59.17) is very close to the median (60.00), indicating",
                    "the distribution is approximately symmetric. The mean is slightly",
                    "less than the median, which can suggest a minor left tail, but the",
                    "difference is negligible\u2014the distribution of physical activity",
                    "levels is effectively symmetric.",
                ]
                clear_and_set_text(shape.text_frame, content, font_size=14)
                break

# ============================================================
# Slide 5 (index 4): Analysis of Daily Steps Taken
# ============================================================
slide_ds = prs.slides[4]
for shape in slide_ds.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '[TODO]' in para.text:
                content = [
                    [("Measures of Spread", True)],
                    "",
                    [("\u2022  Standard Deviation: ", True), ("1,617.92 steps", False)],
                    [("\u2022  Maximum: ", True), ("10,000 steps", False)],
                    [("\u2022  Minimum: ", True), ("3,000 steps", False)],
                    [("\u2022  Range: ", True), ("7,000 steps", False)],
                    "",
                    [("Additional Measures of Spread", True)],
                    "",
                    [("\u2022  Variance: ", True), ("2,617,651.14 steps\u00b2", False)],
                    [("\u2022  Q1 (25th Percentile): ", True), ("5,600 steps", False)],
                    [("\u2022  Q3 (75th Percentile): ", True), ("8,000 steps", False)],
                    [("\u2022  Interquartile Range (IQR): ", True), ("2,400 steps", False)],
                    "",
                    "The IQR of 2,400 steps shows that the middle 50% of participants",
                    "take between 5,600 and 8,000 steps daily.",
                ]
                clear_and_set_text(shape.text_frame, content, font_size=14)
                break

# ============================================================
# Slide 6 (index 5): Distribution of Heart Rates
# ============================================================
slide_hr = prs.slides[5]
for shape in slide_hr.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '[TODO]' in para.text:
                content = [
                    [("Measures of Center: ", True), ("Mean = 70.17,", False)],
                    "Median = 70.00, Mode = 68 bpm",
                    "",
                    [("Distribution Shape: ", True), ("Right-skewed", False)],
                    "Mean (70.17) > Median (70.00)",
                    "confirms a right-skewed distribution.",
                    "",
                    [("Outliers: ", True), ("15 detected (IQR method)", False)],
                    "Values above 78 bpm are outliers:",
                    "80, 81, 82, 83, 84, 85, 86 bpm",
                    "",
                    "Q1=68, Q3=72, IQR=4",
                    "Upper fence = Q3+1.5\u00d7IQR = 78",
                ]
                clear_and_set_text(shape.text_frame, content, font_size=12)
                # Resize text box to left half only
                shape.left = Inches(0.5)
                shape.width = Inches(4.5)
                break

# Insert the heart rate histogram image - right side, no overlap
slide_hr.shapes.add_picture(
    'heart_rate_distribution.png',
    Inches(5.2), Inches(2.0), Inches(5.0), Inches(4.2)
)

# ============================================================
# Save the completed presentation
# ============================================================
output_file = 'Sleep_Health_and_Lifestyle_Analysis.pptx'
prs.save(output_file)
print(f"Presentation saved to: {output_file}")
